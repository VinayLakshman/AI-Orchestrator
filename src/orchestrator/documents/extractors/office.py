from __future__ import annotations

import hashlib
import re
from typing import Any

from ..artifact import DocumentArtifact, DocumentCategory, DocumentPage, DocumentSection, DocumentTable


def _word_count(text: str) -> int:
    return len(re.findall(r"\S+", text))


def _line_count(text: str) -> int:
    if not text:
        return 0
    return text.count("\n") + 1


def _artifact_failure(
    *,
    raw: bytes,
    filename: str,
    mime_type: str,
    content_type: str,
    error: str,
    exception: Exception,
    metadata: dict[str, Any] | None = None,
) -> DocumentArtifact:
    checksum = hashlib.sha256(raw).hexdigest()
    return DocumentArtifact(
        id=checksum[:16],
        filename=filename,
        extension=filename.lower().rpartition(".")[-1] if "." in filename else "",
        mime_type=mime_type,
        size=len(raw),
        checksum=checksum,
        category=DocumentCategory.RICH,
        content_type=content_type,
        metadata={**dict(metadata or {}), "status": "failed", "error": error},
        status="failed",
        error=str(exception),
    )


def extract_docx_document(
    *,
    raw: bytes,
    filename: str,
    mime_type: str,
    metadata: dict[str, Any] | None = None,
) -> DocumentArtifact:
    checksum = hashlib.sha256(raw).hexdigest()
    try:
        from docx import Document  # type: ignore
    except Exception as exc:
        return _artifact_failure(
            raw=raw,
            filename=filename,
            mime_type=mime_type,
            content_type="docx",
            error="python_docx_not_available",
            exception=exc,
            metadata=metadata,
        )

    try:
        doc = Document(None)
        from io import BytesIO

        doc = Document(BytesIO(raw))
        sections: list[DocumentSection] = []
        tables: list[DocumentTable] = []
        parts: list[str] = []
        for paragraph in doc.paragraphs:
            text = paragraph.text or ""
            if text:
                parts.append(text)
        for table_index, table in enumerate(doc.tables, start=1):
            rows: list[list[str]] = []
            row_texts: list[str] = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)
                row_texts.append("\t".join(cells))
            if row_texts:
                parts.append("\n".join(row_texts))
            tables.append(
                DocumentTable(
                    name=f"table_{table_index}",
                    text="\n".join(row_texts),
                    rows=rows,
                    metadata={"table_index": table_index},
                )
            )

        core = getattr(doc, "core_properties", None)
        metadata_map = dict(metadata or {})
        if core is not None:
            metadata_map.update(
                {
                    "author": getattr(core, "author", None),
                    "title": getattr(core, "title", None),
                    "subject": getattr(core, "subject", None),
                    "category": getattr(core, "category", None),
                    "keywords": getattr(core, "keywords", None),
                }
            )

        text = "\n".join(parts).strip()
        sections.append(DocumentSection(title="document", text=text[:4000], metadata={"source": "docx"}))

        return DocumentArtifact(
            id=checksum[:16],
            filename=filename,
            extension="docx",
            mime_type=mime_type,
            size=len(raw),
            checksum=checksum,
            category=DocumentCategory.RICH,
            language="docx",
            encoding="utf-8",
            content_type="docx",
            metadata=metadata_map,
            text=text,
            pages=[DocumentPage(number=1, text=text, metadata={"source": "docx"})] if text else [],
            tables=tables,
            sections=sections,
            line_count=_line_count(text),
            page_count=1 if text else 0,
            word_count=_word_count(text),
            status="ok",
        )
    except Exception as exc:
        return _artifact_failure(
            raw=raw,
            filename=filename,
            mime_type=mime_type,
            content_type="docx",
            error="docx_extract_failed",
            exception=exc,
            metadata=metadata,
        )


def extract_xlsx_document(
    *,
    raw: bytes,
    filename: str,
    mime_type: str,
    metadata: dict[str, Any] | None = None,
) -> DocumentArtifact:
    checksum = hashlib.sha256(raw).hexdigest()
    try:
        from io import BytesIO
        from openpyxl import load_workbook  # type: ignore
    except Exception as exc:
        return _artifact_failure(
            raw=raw,
            filename=filename,
            mime_type=mime_type,
            content_type="xlsx",
            error="openpyxl_not_available",
            exception=exc,
            metadata=metadata,
        )

    try:
        wb = load_workbook(BytesIO(raw), data_only=False)
        parts: list[str] = []
        pages: list[DocumentPage] = []
        for sheet_index, ws in enumerate(wb.worksheets, start=1):
            rows_text: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = ["" if value is None else str(value) for value in row]
                line = "\t".join(cells).rstrip()
                if line.strip():
                    rows_text.append(line)
            sheet_text = "\n".join(rows_text).strip()
            if sheet_text:
                parts.append(f"[Sheet] {ws.title}\n{sheet_text}")
            pages.append(
                DocumentPage(
                    number=sheet_index,
                    text=sheet_text,
                    metadata={"sheet_name": ws.title, "row_count": ws.max_row, "column_count": ws.max_column},
                )
            )

        metadata_map = dict(metadata or {})
        metadata_map.update(
            {
                "sheet_names": wb.sheetnames,
                "active_sheet": wb.active.title if wb.active else None,
            }
        )
        text = "\n\f\n".join(parts).strip()
        return DocumentArtifact(
            id=checksum[:16],
            filename=filename,
            extension="xlsx",
            mime_type=mime_type,
            size=len(raw),
            checksum=checksum,
            category=DocumentCategory.RICH,
            language="spreadsheet",
            encoding="utf-8",
            content_type="xlsx",
            metadata=metadata_map,
            text=text,
            pages=pages,
            line_count=_line_count(text),
            page_count=len(pages),
            word_count=_word_count(text),
            status="ok",
        )
    except Exception as exc:
        return _artifact_failure(
            raw=raw,
            filename=filename,
            mime_type=mime_type,
            content_type="xlsx",
            error="xlsx_extract_failed",
            exception=exc,
            metadata=metadata,
        )


def extract_pptx_document(
    *,
    raw: bytes,
    filename: str,
    mime_type: str,
    metadata: dict[str, Any] | None = None,
) -> DocumentArtifact:
    checksum = hashlib.sha256(raw).hexdigest()
    try:
        from io import BytesIO
        from pptx import Presentation  # type: ignore
    except Exception as exc:
        return _artifact_failure(
            raw=raw,
            filename=filename,
            mime_type=mime_type,
            content_type="pptx",
            error="python_pptx_not_available",
            exception=exc,
            metadata=metadata,
        )

    try:
        prs = Presentation(BytesIO(raw))
        parts: list[str] = []
        pages: list[DocumentPage] = []
        for slide_index, slide in enumerate(prs.slides, start=1):
            slide_lines: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "") or ""
                if text.strip():
                    slide_lines.append(text.strip())
            slide_text = "\n".join(slide_lines).strip()
            if slide_text:
                parts.append(f"[Slide] {slide_index}\n{slide_text}")
            pages.append(
                DocumentPage(
                    number=slide_index,
                    text=slide_text,
                    metadata={"slide_index": slide_index, "shape_count": len(slide.shapes)},
                )
            )

        metadata_map = dict(metadata or {})
        metadata_map.update(
            {
                "slide_count": len(prs.slides),
                "presentation_title": getattr(prs.core_properties, "title", None),
            }
        )
        text = "\n\f\n".join(parts).strip()
        return DocumentArtifact(
            id=checksum[:16],
            filename=filename,
            extension="pptx",
            mime_type=mime_type,
            size=len(raw),
            checksum=checksum,
            category=DocumentCategory.RICH,
            language="presentation",
            encoding="utf-8",
            content_type="pptx",
            metadata=metadata_map,
            text=text,
            pages=pages,
            line_count=_line_count(text),
            page_count=len(pages),
            word_count=_word_count(text),
            status="ok",
        )
    except Exception as exc:
        return _artifact_failure(
            raw=raw,
            filename=filename,
            mime_type=mime_type,
            content_type="pptx",
            error="pptx_extract_failed",
            exception=exc,
            metadata=metadata,
        )

